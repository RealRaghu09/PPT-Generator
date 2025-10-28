import webcolors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime
import os

class ppt:
    def __init__(self):
        # Holds transient data for the current generation request
        self.collected_data = []
    def extract_rgb_values(self, color: str):
        try:
            rgb_value = webcolors.name_to_rgb(color.lower())
            return list(rgb_value)
        except Exception:
            return [0, 0, 0]  # Default to black

    def add_thankyou_slide(self, prs, font, font_color, bg_color):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Set background color
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*bg_color)

        # Add centered "Thank You"
        box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        tf = box.text_frame
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Thank You"
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.name = font
        run.font.color.rgb = RGBColor(*font_color)

    def generate_normal_ppt(self, topic: str, list_of_content: list[str], list_of_topics: list[str], no_of_slides: int,
                     color_scheme: list[str], fontstyle: list[str]) -> str:
        # Collect input data for this generation
        self.collected_data.append({
            "style": "normal",
            "topic": topic,
            "list_of_topics": list_of_topics,
            "list_of_content": list_of_content,
            "no_of_slides": no_of_slides,
            "color_scheme": color_scheme,
            "fontstyle": fontstyle,
        })
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        rgb1, rgb2 = [0, 0, 0], [255, 255, 255]
        if len(color_scheme) >= 2:
            rgb1 = self.extract_rgb_values(color_scheme[0])
            rgb2 = self.extract_rgb_values(color_scheme[1])

        title_font = fontstyle[0] if len(fontstyle) > 0 else "Arial"
        content_font = fontstyle[1] if len(fontstyle) > 1 else "Calibri"

        # First slide
        first_slide = prs.slides.add_slide(blank_slide_layout)
        title_shape = first_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        tf = title_shape.text_frame
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = topic
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.name = title_font
        run.font.color.rgb = RGBColor(*rgb1)
        first_slide.background.fill.solid()
        first_slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

        # Slides 2 to no_of_slides
        for i in range(min(no_of_slides - 1, len(list_of_topics), len(list_of_content))):
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

            # Title
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1))
            title_tf = title_box.text_frame
            title_tf.auto_size = True
            p_title = title_tf.paragraphs[0]
            p_title.alignment = PP_ALIGN.CENTER
            run_title = p_title.add_run()
            run_title.text = list_of_topics[i]
            run_title.font.size = Pt(16)
            run_title.font.bold = True
            run_title.font.name = title_font
            run_title.font.color.rgb = RGBColor(*rgb1)

            # Content
            # New line (adds more vertical gap between title and content)
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(5))
            content_tf = content_box.text_frame
            content_tf.auto_size = True
            p_content = content_tf.paragraphs[0]
            p_content.alignment = PP_ALIGN.CENTER
            run_content = p_content.add_run()
            run_content.text = list_of_content[i]
            run_content.font.size = Pt(14)
            run_content.font.name = content_font
            run_content.font.color.rgb = RGBColor(*rgb1)

        # ➕ Adding "Thank You" slide at end
        self.add_thankyou_slide(prs, font=title_font, font_color=tuple(rgb1), bg_color=tuple(rgb2))

        # Save with hashed filename
        timestamp = datetime.datetime.now().isoformat()
        unique_hash = (topic + timestamp)
        filename = f"generated_ppt/generated_presentation_{unique_hash}.pptx"

        os.makedirs("generated_ppt", exist_ok=True)
        prs.save(filename)

        # Clear transient data after generation completes
        self.collected_data.clear()

        return {"message": f"Generated Normal Successfully! Saved as {filename}", "file_path": filename}
    

    def generate_modern_ppt(self, topic: str, list_of_content: list[str], list_of_topics: list[str], no_of_slides: int) -> str:
        # Collect input data for this generation
        self.collected_data.append({
            "style": "modern",
            "topic": topic,
            "list_of_topics": list_of_topics,
            "list_of_content": list_of_content,
            "no_of_slides": no_of_slides,
        })

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        # Color palette: title color (accent), bg color (light/dark)
        rgb1 = [30, 30, 30]   # Dark Gray/Black
        rgb2 = [245, 245, 245] # Light Gray

        title_font = "Montserrat"
        content_font =  "Lato"

        # First slide: Title
        first_slide = prs.slides.add_slide(blank_slide_layout)
        first_slide.background.fill.solid()
        first_slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

        title_box = first_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        tf = title_box.text_frame
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = topic
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.name = title_font
        run.font.color.rgb = RGBColor(*rgb1)

        # Content slides
        for i in range(min(no_of_slides - 1, len(list_of_topics), len(list_of_content))):
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

            # Title
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(1))
            title_tf = title_box.text_frame
            title_tf.auto_size = True
            p_title = title_tf.paragraphs[0]
            p_title.alignment = PP_ALIGN.LEFT
            run_title = p_title.add_run()
            run_title.text = list_of_topics[i]
            run_title.font.size = Pt(28)
            run_title.font.bold = True
            run_title.font.name = title_font
            run_title.font.color.rgb = RGBColor(*rgb1)

            # Content
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8.5), Inches(5))
            content_tf = content_box.text_frame
            content_tf.word_wrap = True
            content_tf.auto_size = True
            p_content = content_tf.paragraphs[0]
            p_content.alignment = PP_ALIGN.LEFT
            run_content = p_content.add_run()
            run_content.text = list_of_content[i]
            run_content.font.size = Pt(18)
            run_content.font.name = content_font
            run_content.font.color.rgb = RGBColor(*rgb1)

        # Add modern-style thank you slide
        self.add_thankyou_slide(prs, font=title_font, font_color=tuple(rgb1), bg_color=tuple(rgb2))

        # Save presentation
        timestamp = datetime.datetime.now().isoformat()
        unique_hash = (topic + timestamp)
        filename = f"generated_ppt/modern_presentation_{unique_hash}.pptx"
        os.makedirs("generated_ppt", exist_ok=True)
        prs.save(filename)
        
        # Clear transient data after generation completes
        self.collected_data.clear()

        return {"message":f"Modern PPT Generated Successfully! Saved as {filename}" , "file_path" : filename}

    def generate_creative_ppt(self, topic: str, list_of_content: list[str], list_of_topics: list[str], no_of_slides: int) -> str:
        # Collect input data for this generation
        self.collected_data.append({
            "style": "creative",
            "topic": topic,
            "list_of_topics": list_of_topics,
            "list_of_content": list_of_content,
            "no_of_slides": no_of_slides,
        })
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        import datetime, os

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        # Colors
        rgb1 =  [255, 255, 255]  # Text
        rgb2 =  [0, 0, 0]        # Background

        title_font =  "Poppins"
        content_font = "Comic Sans MS"

        # First slide (Title slide - center aligned and bold)
        first_slide = prs.slides.add_slide(blank_slide_layout)
        first_slide.background.fill.solid()
        first_slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

        title_box = first_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2.5))
        tf = title_box.text_frame
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = topic.upper()
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.name = title_font
        run.font.color.rgb = RGBColor(*rgb1)

        # Content slides
        for i in range(min(no_of_slides - 1, len(list_of_topics), len(list_of_content))):
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

            # Title (bold and colored)
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.5), Inches(1))
            title_tf = title_box.text_frame
            title_tf.auto_size = True
            p_title = title_tf.paragraphs[0]
            p_title.alignment = PP_ALIGN.LEFT
            run_title = p_title.add_run()
            run_title.text = f"🎨 {list_of_topics[i]}"
            run_title.font.size = Pt(32)
            run_title.font.bold = True
            run_title.font.name = title_font
            run_title.font.color.rgb = RGBColor(*rgb1)

            # Content (centered and expressive)
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8.5), Inches(5))
            content_tf = content_box.text_frame
            content_tf.word_wrap = True
            content_tf.auto_size = True
            p_content = content_tf.paragraphs[0]
            p_content.alignment = PP_ALIGN.LEFT
            run_content = p_content.add_run()
            run_content.text = list_of_content[i]
            run_content.font.size = Pt(20)
            run_content.font.name = content_font
            run_content.font.color.rgb = RGBColor(*rgb1)

        # Add creative thank-you slide
        self.add_thankyou_slide(prs, font=title_font, font_color=tuple(rgb1), bg_color=tuple(rgb2))

        # Save file
        timestamp = datetime.datetime.now().isoformat()
        unique_hash = hash(topic + timestamp) % (10 ** 8)
        filename = f"generated_ppt/creative_presentation_{unique_hash}.pptx"
        os.makedirs("generated_ppt", exist_ok=True)
        prs.save(filename)
        
        # Clear transient data after generation completes
        self.collected_data.clear()

        return {"message" : f"Creative PPT Generated Successfully! Saved as {filename}" , "file_path" : filename}
    def generate_retro_ppt(self, topic: str, list_of_content: list[str], list_of_topics: list[str], no_of_slides: int) -> str:
        # Collect input data for this generation
        self.collected_data.append({
            "style": "retro",
            "topic": topic,
            "list_of_topics": list_of_topics,
            "list_of_content": list_of_content,
            "no_of_slides": no_of_slides,
        })

        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]

        # Retro default palette (fallbacks)
        rgb1 = [139, 69, 19]   # Saddle Brown
        rgb2 = [255, 228, 181] # Moccasin

        title_font = "Courier New"
        content_font = "Georgia"

        # First slide - Title
        first_slide = prs.slides.add_slide(blank_slide_layout)
        first_slide.background.fill.solid()
        first_slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

        title_box = first_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2.5))
        tf = title_box.text_frame
        tf.auto_size = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = topic.upper()
        run.font.size = Pt(42)
        run.font.bold = True
        run.font.name = title_font
        run.font.color.rgb = RGBColor(*rgb1)

        # Content slides
        for i in range(min(no_of_slides - 1, len(list_of_topics), len(list_of_content))):
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*rgb2)

            # Title
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(1))
            title_tf = title_box.text_frame
            title_tf.auto_size = True
            p_title = title_tf.paragraphs[0]
            p_title.alignment = PP_ALIGN.LEFT
            run_title = p_title.add_run()
            run_title.text = f"★ {list_of_topics[i]}"
            run_title.font.size = Pt(28)
            run_title.font.bold = True
            run_title.font.name = title_font
            run_title.font.color.rgb = RGBColor(*rgb1)

            # Content
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8.5), Inches(5))
            content_tf = content_box.text_frame
            content_tf.word_wrap = True
            content_tf.auto_size = True
            p_content = content_tf.paragraphs[0]
            p_content.alignment = PP_ALIGN.LEFT
            run_content = p_content.add_run()
            run_content.text = list_of_content[i]
            run_content.font.size = Pt(18)
            run_content.font.name = content_font
            run_content.font.color.rgb = RGBColor(*rgb1)

        # Thank You Slide
        self.add_thankyou_slide(prs, font=title_font, font_color=tuple(rgb1), bg_color=tuple(rgb2))

        # Save
        timestamp = datetime.datetime.now().isoformat()
        unique_hash = hash(topic + timestamp) % (10 ** 8)
        filename = f"generated_ppt/retro_presentation_{unique_hash}.pptx"
        os.makedirs("generated_ppt", exist_ok=True)
        prs.save(filename)
        
        # Clear transient data after generation completes
        self.collected_data.clear()

        return {"message":f"Retro PPT Generated Successfully! Saved as {filename}" , "file_path" : filename}
